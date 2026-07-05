"""
Document processing and RAG pipeline.
"""

import uuid
from pathlib import Path
from typing import Optional

from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

from app.core.config import get_settings

settings = get_settings()


class DocumentProcessor:
    """Service for processing documents and indexing to Qdrant."""

    def __init__(self):
        """Initialize Qdrant client."""
        self.qdrant_client = QdrantClient(
            host=settings.qdrant_host,
            port=settings.qdrant_port,
        )
        self.chunk_size = settings.rag_chunk_size
        self.chunk_overlap = settings.rag_chunk_overlap
        # Ensure collection exists on initialization
        self._ensure_collection_exists()

    def _extract_text_from_file(self, file_path: str, file_type: str) -> str:
        """
        Extract text from a file based on its type.

        Args:
            file_path: Path to the file
            file_type: File extension/type

        Returns:
            Extracted text
        """
        file_path = Path(file_path)

        if file_type.lower() == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif file_type.lower() == "md":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif file_type.lower() == "html":
            try:
                from html.parser import HTMLParser

                class MLStripper(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.reset()
                        self.strict = False
                        self.convert_charrefs = True
                        self.text = []

                    def handle_data(self, d):
                        self.text.append(d)

                    def get_data(self):
                        return "".join(self.text)

                with open(file_path, "r", encoding="utf-8") as f:
                    stripper = MLStripper()
                    stripper.feed(f.read())
                    return stripper.get_data()
            except Exception as e:
                print(f"Error extracting HTML: {e}")
                return ""

        elif file_type.lower() == "pdf":
            try:
                import PyPDF2

                text = []
                with open(file_path, "rb") as f:
                    try:
                        reader = PyPDF2.PdfReader(f)
                        for page in reader.pages:
                            try:
                                extracted = page.extract_text()
                                if extracted:
                                    text.append(extracted)
                            except Exception:
                                pass
                    except Exception as e:
                        print(f"PyPDF2 error: {e}")
                        # Fallback: try pypdf (alternative name)
                        try:
                            import pypdf
                            f.seek(0)
                            reader = pypdf.PdfReader(f)
                            for page in reader.pages:
                                try:
                                    extracted = page.extract_text()
                                    if extracted:
                                        text.append(extracted)
                                except Exception:
                                    pass
                        except Exception:
                            pass
                
                result = "\n".join(text).strip()
                if result:
                    return result
                else:
                    # If text extraction failed, return placeholder
                    return f"[PDF Document: {file_path} - Text extraction returned empty]\n"
            except Exception as e:
                print(f"Error extracting PDF: {e}")
                return f"[PDF Document: {file_path} - Error: {str(e)}]\n"

        elif file_type.lower() == "docx":
            try:
                from docx import Document

                doc = Document(file_path)
                text = []
                for para in doc.paragraphs:
                    text.append(para.text)
                return "\n".join(text)
            except Exception as e:
                print(f"Error extracting DOCX: {e}")
                return ""

        elif file_type.lower() == "pptx":
            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            except Exception as e:
                print(f"Error extracting PPTX: {e}")
                return ""

        return ""

    def _chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> list[str]:
        """
        Split text into overlapping chunks.

        Args:
            text: Text to chunk
            chunk_size: Size of each chunk (default from config)
            overlap: Number of characters to overlap (default from config)

        Returns:
            List of text chunks
        """
        chunk_size = chunk_size or self.chunk_size
        overlap = overlap or self.chunk_overlap

        chunks = []
        start = 0

        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start = end - overlap

        return [chunk.strip() for chunk in chunks if chunk.strip()]

    def _simple_hash(self, text: str) -> int:
        """
        Create a simple hash of text for use as vector ID.

        Args:
            text: Text to hash

        Returns:
            Hash value
        """
        return hash(text) & 0x7FFFFFFF

    def process_document(
        self,
        document_id: uuid.UUID,
        file_path: str,
        file_type: str,
        filename: str,
    ) -> tuple[int, Optional[str]]:
        """
        Process a document and index it to Qdrant.

        Args:
            document_id: ID of the document
            file_path: Path to the file
            file_type: File extension
            filename: Original filename

        Returns:
            Tuple of (chunk_count, error_message)
            - chunk_count: Number of chunks created (0 if error)
            - error_message: Error message if failed, None if success
        """
        try:
            # Extract text from file
            text = self._extract_text_from_file(file_path, file_type)
            if not text or len(text.strip()) == 0:
                return 0, f"No text could be extracted from {filename}"

            # Chunk the text
            chunks = self._chunk_text(text)
            if not chunks:
                return 0, f"No chunks generated from {filename}"

            # Ensure collection exists
            self._ensure_collection_exists()

            # Index chunks to Qdrant
            points = []
            for i, chunk in enumerate(chunks):
                point_id = self._simple_hash(f"{document_id}_{i}_{chunk[:50]}")
                points.append(
                    PointStruct(
                        id=point_id,
                        vector=[0.0] * 384,  # Dummy vector for text-based search
                        payload={
                            "document_id": str(document_id),
                            "chunk_index": i,
                            "text": chunk,
                            "filename": filename,
                            "file_type": file_type,
                            "metadata": {
                                "chunk_index": i,
                                "total_chunks": len(chunks),
                            },
                        },
                    )
                )

            # Upload points to Qdrant
            self.qdrant_client.upsert(
                collection_name=settings.qdrant_collection_name,
                points=points,
            )

            return len(chunks), None

        except Exception as e:
            error_msg = f"Error processing document: {str(e)}"
            print(error_msg)
            return 0, error_msg

    def _ensure_collection_exists(self):
        """Ensure Qdrant collection exists, create if not."""
        try:
            self.qdrant_client.get_collection(settings.qdrant_collection_name)
        except Exception:
            # Collection doesn't exist, create it
            self.qdrant_client.create_collection(
                collection_name=settings.qdrant_collection_name,
                vectors_config=VectorParams(
                    size=384,
                    distance=Distance.COSINE,
                ),
            )

    def delete_document_chunks(self, document_id: uuid.UUID) -> bool:
        """
        Delete all chunks for a document from Qdrant.

        Args:
            document_id: ID of the document to delete

        Returns:
            True if successful, False otherwise
        """
        try:
            self.qdrant_client.delete(
                collection_name=settings.qdrant_collection_name,
                points_selector={
                    "filter": {
                        "must": [
                            {
                                "key": "document_id",
                                "match": {"value": str(document_id)},
                            }
                        ]
                    }
                },
            )
            return True
        except Exception as e:
            print(f"Error deleting document chunks: {e}")
            return False
